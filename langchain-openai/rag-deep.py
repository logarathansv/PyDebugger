import os
import streamlit as st
from langchain.schema import Document
from langchain_community.document_loaders import PDFPlumberLoader
from langchain_text_splitters import MarkdownTextSplitter, RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_core.prompts import ChatPromptTemplate
from langchain_community.embeddings import OllamaEmbeddings
from langchain_ollama import OllamaLLM
from youtube_transcript_api import YouTubeTranscriptApi
from urllib.parse import urlparse, parse_qs
import json
import time
from pptx import Presentation
import ast
import sys
import traceback
from io import StringIO
from contextlib import redirect_stdout, redirect_stderr
import re
# Initialize with deepseek-r1:1.5b model
OLLAMA_MODEL = "deepseek-r1:1.5b"

# Initialize LLM and Embeddings
llm = OllamaLLM(model=OLLAMA_MODEL, temperature=0.3)
embeddings = OllamaEmbeddings(model=OLLAMA_MODEL)

# UI Configuration
st.set_page_config(page_title="Learning Chatbot & Debugging Assistantt", layout="wide")
st.title("💡Learning Chatbot & Debugging Assistant")

# File handling
PDF_STORAGE_PATH = 'document_store/'
CHROMA_DB_PATH = 'chroma_db/'
os.makedirs(PDF_STORAGE_PATH, exist_ok=True)
os.makedirs(CHROMA_DB_PATH, exist_ok=True)

# Dark Mode CSS
st.markdown("""
    <style>
    .stApp { background-color: #0E1117; color: #FFFFFF; }
    .stChatInput input { background-color: #1E1E1E !important; color: #FFFFFF !important; border: 1px solid #3A3A3A !important; }
    .stChatMessage p, .stChatMessage div { color: #FFFFFF !important; }
    .stFileUploader { background-color: #1E1E1E; border: 1px solid #3A3A3A; border-radius: 5px; padding: 15px; }
    h1, h2, h3 { color: #00FFAA !important; }
    </style>
    """, unsafe_allow_html=True)

def save_uploaded_file(uploaded_file):
    file_path = os.path.join(PDF_STORAGE_PATH, uploaded_file.name)
    with open(file_path, "wb") as f:
        f.write(uploaded_file.getbuffer())
    return file_path

def extract_text_from_pptx(file_path):
    presentation = Presentation(file_path)
    return "\n".join([shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")])

def process_document(file_path, file_name):
    file_ext = file_name.split(".")[-1].lower()
    
    if file_ext == "pdf":
        loader = PDFPlumberLoader(file_path)
        docs = loader.load()
    elif file_ext == "txt":
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
        docs = [Document(page_content=text)]
    elif file_ext == "pptx":
        text = extract_text_from_pptx(file_path)
        docs = [Document(page_content=text)]
    else:
        st.error("Unsupported file format")
        return None
    
    # Choose splitter based on file type
    if file_ext in ["md", "rst"]:
        splitter = MarkdownTextSplitter(chunk_size=500, chunk_overlap=100)
    else:
        splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=100)
    
    chunks = splitter.split_documents(docs)
    
    # Create vector store
    collection_name = file_name.replace(" ", "_").lower()
    vector_store = Chroma.from_documents(
        documents=chunks,
        embedding=embeddings,
        collection_name=collection_name,
        persist_directory=CHROMA_DB_PATH
    )
    
    return vector_store

def extract_youtube_transcript(video_url):
    try:
        video_id = parse_qs(urlparse(video_url).query).get("v", [None])[0]
        if not video_id:
            return None, "Invalid YouTube URL"
            
        transcript = YouTubeTranscriptApi.get_transcript(video_id, languages=['en'])
        text = " ".join([entry['text'] for entry in transcript])
        
        # Save transcript
        file_name = f"youtube_{video_id}.txt"
        file_path = os.path.join(PDF_STORAGE_PATH, file_name)
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(text)
            
        return file_path, None
    except Exception as e:
        return None, str(e)

def find_context(query, selected_docs, mode):
    context_docs = []
    citations = []

    for doc_name in selected_docs:
        vector_store = st.session_state.vector_stores.get(doc_name)
        if not vector_store:
            continue
        
        results = vector_store.similarity_search_with_score(query, k=2)
        
        for doc, score in results:
            if score > 0.7:  # Only include relevant results
                context_docs.append(doc)
                citations.append((doc_name, doc.metadata.get("page", "N/A")))
    
    return context_docs, citations

def generate_follow_up_hint(chat_history, mode):
    conversation_context = "\n".join([f"{msg['role']}: {msg['content']}" for msg in chat_history])
    
    if mode == "Programming Tutor":
        prompt = """You are a programming tutor. Provide another hint based on this conversation.
        Keep it concise and helpful.

        Conversation:
        {conversation_context}

        Hint:"""
    else:
        prompt = """You are a debugging assistant. Provide another guiding question.

        Conversation:
        {conversation_context}

        Hint:"""

    return llm(prompt.format(conversation_context=conversation_context))

def generate_answer(query, context_docs, citations, mode):
    context = "\n\n".join([doc.page_content[:500] for doc in context_docs])
    
    if mode == "Programming Tutor":
        prompt = """You are a programming tutor. Explain concepts clearly with examples.
        Use this context:
        {context}
        
        Question: {query}
        Answer:"""
    else:
        prompt = """You are a debugging assistant. Provide step-by-step guidance.
        Context:
        {context}
        
        Problem: {query}
        Guidance:"""

    answer = llm(prompt.format(context=context, query=query))
    
    # Add citations
    if citations:
        sources = "\n".join([f"- {pdf}" for pdf, _ in citations])
        answer += f"\n\nSources:\n{sources}"
    
    return answer

# Initialize session state
if "vector_stores" not in st.session_state:
    st.session_state.vector_stores = {}
    
if "messages" not in st.session_state:
    st.session_state.messages = []

# Sidebar
with st.sidebar:
    st.header("📂 Document Management")
    
    # File upload
    uploaded_files = st.file_uploader("Upload documents", 
                                    type=["pdf", "txt", "pptx", "md"], 
                                    accept_multiple_files=True)
    
    if uploaded_files:
        for file in uploaded_files:
            if file.name not in st.session_state.vector_stores:
                file_path = save_uploaded_file(file)
                vector_store = process_document(file_path, file.name)
                if vector_store:
                    st.session_state.vector_stores[file.name] = vector_store
                    st.success(f"Processed {file.name}")
    
    # YouTube transcript
    st.header("🎬 YouTube Transcript")
    youtube_url = st.text_input("Enter YouTube URL")
    if st.button("Get Transcript") and youtube_url:
        file_path, error = extract_youtube_transcript(youtube_url)
        if error:
            st.error(f"Error: {error}")
        else:
            vector_store = process_document(file_path, f"youtube_{os.path.basename(file_path)}")
            if vector_store:
                st.session_state.vector_stores[f"youtube_{os.path.basename(file_path)}"] = vector_store
                st.success("Transcript processed!")
    
    # Document selection
    st.header("📑 Active Documents")
    selected_docs = []
    for doc_name in st.session_state.vector_stores.keys():
        if st.checkbox(doc_name, value=True):
            selected_docs.append(doc_name)

    # Assistant mode selection
    mode = st.radio("Assistant Mode:", ["Programming Tutor", "Debugging Assistant"])

# Chat interface
for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])

if prompt := st.chat_input("Ask a question..."):
    st.session_state.messages.append({"role": "user", "content": prompt})
    with st.chat_message("user"):
        st.markdown(prompt)
    
    # Generate response
    with st.chat_message("assistant"):
        with st.spinner("Thinking..."):
            try:
                context_docs, citations = find_context(prompt, selected_docs, mode)
                answer = generate_answer(prompt, context_docs, citations, mode)
                think_match = re.search(r"<think>(.*?)</think>", answer, re.DOTALL)
                think_text = think_match.group(1).strip() if think_match else ""

# Remove <think>...</think> from the main message
                cleaned_answer = re.sub(r"<think>.*?</think>", "", answer, flags=re.DOTALL).strip()

                st.markdown(answer)
                st.session_state.messages.append({"role": "assistant", "content": cleaned_answer})
            except Exception as e:
                st.error(f"Error: {str(e)}")

# Follow-up hint button
if mode == "Debugging Assistant" and st.session_state.messages and st.session_state.messages[-1]["role"] == "assistant":
    if st.button("💡 Need another hint"):
        with st.spinner("Generating hint..."):
            try:
                hint = generate_follow_up_hint(st.session_state.messages, mode)
                think_match = re.search(r"<think>(.*?)</think>", hint, re.DOTALL)
                think_text = think_match.group(1).strip() if think_match else ""

# Remove <think>...</think> from the main message
                cleaned_hint = re.sub(r"<think>.*?</think>", "", hint, flags=re.DOTALL).strip()

                st.session_state.messages.append({"role": "assistant", "content": cleaned_hint})
                st.rerun()
            except Exception as e:
                st.error(f"Error generating hint: {str(e)}")

#-------------------------------------------------------------------------------------------
# Add this to your existing imports

# Add this class definition somewhere before your main UI code
class DebugSandbox:
    def __init__(self):
        self.breakpoints = set()
        self.current_line = 0
        self.variables = {}
        self.execution_log = []

    def execute_with_debug(self, code):
        """Execute code with debug capabilities"""
        tree = ast.parse(code)
        self._instrument_code(tree)

        output = StringIO()
        error = StringIO()

        # Add self reference to the execution namespace
        self.variables['self'] = self

        try:
            ast.fix_missing_locations(tree)
            with redirect_stdout(output), redirect_stderr(error):
                exec(compile(tree, '<string>', 'exec'), self.variables)
        except Exception:
            error.write(traceback.format_exc())
        finally:
            # Remove self reference to avoid pollution
            if 'self' in self.variables:
                del self.variables['self']

        return {
            'output': output.getvalue(),
            'error': error.getvalue(),
            'variables': {k: v for k, v in self.variables.items() if k != 'self'},
            'execution_log': self.execution_log
        }

    def _instrument_code(self, node):
        """Add debugging instrumentation to AST"""
        for field, old_value in ast.iter_fields(node):
            if isinstance(old_value, list):
                new_values = []
                for value in old_value:
                    if isinstance(value, ast.stmt):
                        value = self._instrument_statement(value)
                    new_values.append(value)
                setattr(node, field, new_values)

    def _instrument_statement(self, node):
        self.current_line += 1
        log_expr = ast.Expr(value=ast.Call(
            func=ast.Attribute(value=ast.Name(id='self', ctx=ast.Load()), 
            attr='_log_execution', ctx=ast.Load()),
            args=[ast.Constant(value=self.current_line)],
            keywords=[]
        ))

        # Set location info from original node
        log_expr = ast.copy_location(log_expr, node)

        if isinstance(node, ast.FunctionDef):
            node.body = [log_expr] + node.body
            return node
        else:
            new_node = ast.If(
                test=ast.Constant(value=True),
                body=[log_expr, node],
                orelse=[]
            )
            new_node = ast.copy_location(new_node, node)
            return new_node

    def _log_execution(self, line_num):
        """Log execution details"""
        self.execution_log.append({
            'line': line_num,
            'variables': {k: v for k, v in self.variables.items() if k != 'self'},
            'call_stack': traceback.extract_stack()[:-2]
        })
def generate_answer1(query, code):
    # Enhanced debugging prompt
    prompt = """You are a debugging assistant. Analyze the problem and provide step-by-step guidance.
Context:
{code}
Problem: {query}
Provide:
1. Error analysis
2. Step-by-step solution approach
3. Suggested code fixes (if applicable)
Guidance:"""

    # Assuming `llm` is a function that takes a prompt and returns a response from a language model
    answer = llm(prompt.format(code=code, query=query))

    return answer
def generate_follow_up_hint1(chat_history):
    # Extract the conversation context
    conversation_context = "\n".join([f"{msg['role']}: {msg['content']}" for msg in chat_history])
    print(conversation_context)  # Optional: remove or log only in debug mode

    # Define the prompt with correct formatting
    prompt = """You are a rubber duck debugging assistant. Provide another guiding question or hint to help the user think through their problem.

Conversation:
{conversation_context}

Hint:"""

    # Format prompt using ChatPromptTemplate
    conversation_prompt = ChatPromptTemplate.from_template(prompt)

    return llm(prompt.format(conversation_context=conversation_context))


def get_line_code(full_code, line_num):
    """Returns the specific line of code with line number, plus 2 lines of context"""
    lines = full_code.split('\n')
    start_line = max(1, line_num - 1)
    end_line = min(len(lines), line_num + 1)
    
    # Format with line numbers
    result = []
    for i in range(start_line, end_line + 1):
        line_content = lines[i-1]
        prefix = ">>" if i == line_num else f"{i:4}"
        result.append(f"{prefix}: {line_content}")
    
    return '\n'.join(result)

# Modify your sidebar to include the debug sandbox option
with st.sidebar:
    st.header("🔧 Tools")
    if st.checkbox("Show Debug Sandbox", key="show_debug_sandbox"):
        st.session_state.show_sandbox = True
    else:
        st.session_state.show_sandbox = False

# Add this after your main chat interface
if st.session_state.get('show_sandbox', False):
    st.markdown("---")
    st.title("🐞 Python Debugging Sandbox")
    
    # Add some CSS for the debug sandbox
    st.markdown("""
        <style>
        .debug-container {
            border: 1px solid #3A3A3A;
            border-radius: 5px;
            padding: 15px;
            margin-bottom: 20px;
            background-color: #1E1E1E;
        }
        .debug-tabs .stTab {
            background-color: #1E1E1E !important;
        }
        </style>
    """, unsafe_allow_html=True)
    
    code = st.text_area("Enter Python code to debug:", height=300, key="debug_code")
    
    col1, col2 = st.columns(2)
    with col1:
        if st.button("Run with Debugging"):
            sandbox = DebugSandbox()
            result = sandbox.execute_with_debug(code)
            st.session_state.debug_result = result
            st.session_state.stored_code = code
        if st.button("💡 Need another hint"):
            with st.spinner("Generating hint..."):
                try:
                    hint = generate_follow_up_hint1(st.session_state.messages)
                    st.markdown(f"**hint** `{hint}`")
            
                    st.rerun()
                except Exception as e:
                    st.error(f"Error generating hint: {str(e)}")
    
    if 'debug_result' in st.session_state:
        result = st.session_state.debug_result
        code = st.session_state.stored_code
        
        st.markdown("### Execution Results")
        tab1, tab2, tab3 = st.tabs(["Output", "Variables", "Errors"])
        
        with tab1:
            st.code(result['output'] or "No output", language='text')
        
        with tab2:
            st.json(result['variables'])
        
        with tab3:
            if result['error']:
                st.error("**Error Details**")
                st.code(result['error'], language='python')
                
                error_lines = result['error'].strip().split('\n')
                if error_lines:
                    last_line = error_lines[-1]
                    if ':' in last_line:
                        error_type, error_message = last_line.split(':', 1)
                        error_type = error_type.strip()
                        error_message = error_message.strip()
                    else:
                        error_type = "Error"
                        error_message = last_line.strip()
                        answer=generate_answer(error_message,code)
                    st.markdown("**Debug Summary**")
                    st.error(f"**Error Type:** `{error_type}`")
                    st.error(f"**Error Message:** `{error_message}`")
                    st.markdown(f"**Debug Summary from llm** `{generate_answer1(error_message,code)}`")
                   
                    error_line = None
                    for line in error_lines:
                        if "line " in line.lower() and ", in " in line.lower():
                            try:
                                line_part = line.split("line ")[1]
                                line_num = line_part.split(",")[0] if "," in line_part else line_part
                                error_line = int(line_num.strip())
                            except (IndexError, ValueError):
                                continue
                    
                    if error_line is not None:
                        st.error(f"**Error Location:** Line {error_line}")
                        st.code(get_line_code(code, error_line), language='python')
            else:
                st.success("No errors detected")

# Modify your generate_answer function to handle debugging queries
