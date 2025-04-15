import os
from langchain.schema import Document
from openai import BadRequestError
import streamlit as st
from langchain_community.document_loaders import PDFPlumberLoader
from langchain_text_splitters import MarkdownTextSplitter, RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_core.prompts import ChatPromptTemplate
from dotenv import load_dotenv
from langchain_openai import AzureChatOpenAI, AzureOpenAIEmbeddings
import time
import json
from youtube_transcript_api import YouTubeTranscriptApi
from pptx import Presentation
from urllib.parse import urlparse, parse_qs
from sentence_transformers import SentenceTransformer
from langchain.embeddings.base import Embeddings
from typing import List
from together import Together
from docutils.core import publish_parts
from transformers import AutoTokenizer

if "llm_summary_generated" not in st.session_state:
    st.session_state.llm_summary_generated = False

class CustomEmbeddings(Embeddings):
    def __init__(self, model_name: str):
        self.model = SentenceTransformer(model_name)

    def embed_documents(self, documents: List[str]) -> List[List[float]]:
        return [self.model.encode(d).tolist() for d in documents]

    def embed_query(self, query: str) -> List[float]:
        return self.model.encode([query])[0].tolist()

# Load environment variables
load_dotenv()

# Azure OpenAI Embeddings Configuration
AZURE_EMBEDDING_ENDPOINT = os.getenv("AZURE_EMBEDDING_OPENAI")
AZURE_EMBEDDING_API = os.getenv("AZURE_EMBEDDING_API_KEY")
AZURE_EMBEDDING_DEPLOYMENT = os.getenv("AZURE_EMBEDDING_DEPLOYMENT")

EMBEDDING_MODEL = AzureOpenAIEmbeddings(
    model=AZURE_EMBEDDING_DEPLOYMENT,
    azure_endpoint=AZURE_EMBEDDING_ENDPOINT,
    openai_api_version="2023-05-15",
    api_key=AZURE_EMBEDDING_API,
)

analyser = Together(
    api_key="tgp_v1_VlOsAsjxF5mGEzlEBXcGwtVeW25r3gn1OAyci8ZjZJk",
)

EMBEDDING_MODEL_2 = CustomEmbeddings("all-MiniLM-L6-v2")

# Azure OpenAI API Configuration
AZURE_OPENAI_ENDPOINT = os.getenv("AZURE_OPENAI_ENDPOINT")
AZURE_OPENAI_KEY = os.getenv("AZURE_OPENAI_API_KEY")
AZURE_DEPLOYMENT_NAME = os.getenv("AZURE_DEPLOYMENT_NAME")

# Initialize Azure LLM
LANGUAGE_MODEL = AzureChatOpenAI(
    azure_endpoint=AZURE_OPENAI_ENDPOINT,
    openai_api_version="2025-01-01-preview",
    deployment_name=AZURE_DEPLOYMENT_NAME,
    openai_api_key=AZURE_OPENAI_KEY,
    temperature=0.3,
    max_retries=10
)

st.set_page_config(
    page_title="PyDebug",  # Replace with your desired title
    page_icon="🔧")
# Dark Mode UI
st.markdown("""
    <style>
    .stApp { background-color: #0E1117; color: #FFFFFF; }
    .stChatInput input { background-color: #1E1E1E !important; color: #FFFFFF !important; border: 1px solid #3A3A3A !important; }
    .stChatMessage p, .stChatMessage div { color: #FFFFFF !important; }
    .stFileUploader { background-color: #1E1E1E; border: 1px solid #3A3A3A; border-radius: 5px; padding: 15px; }
    h1, h2, h3 { color: #00FFAA !important; }
    </style>
    """, unsafe_allow_html=True)

# System Config
PDF_STORAGE_PATH = 'document_store/pdfs/'
CHROMA_DB_PATH = 'chroma_db/'
CHROMA_DB_PATH_ERROR = 'chroma_db_ERROR/'
PDF_LIST_PATH = os.path.join(CHROMA_DB_PATH, "pdf_list.json")

# Ensure directories exist
os.makedirs(PDF_STORAGE_PATH, exist_ok=True)
os.makedirs(CHROMA_DB_PATH, exist_ok=True)

def format_time(seconds):
    hours = int(seconds // 3600)
    minutes = int((seconds % 3600) // 60)
    secs = int(seconds % 60)
    return f"{hours:02d}:{minutes:02d}:{secs:02d}"

# Function to save transcript
def save_transcript(video_name, transcript_text):
    file_path = os.path.join(PDF_STORAGE_PATH, f"{video_name}.txt")
    with open(file_path, "w", encoding="utf-8") as file:
        file.write(transcript_text)
    return file_path
# Function to extract YouTube Video ID from URL
def extract_video_id(url):
    parsed_url = urlparse(url)
    video_id = None
    
    if parsed_url.netloc in ["www.youtube.com", "youtube.com", "m.youtube.com"]:
        video_id = parse_qs(parsed_url.query).get("v", [None])[0]
    elif parsed_url.netloc in ["youtu.be"]:  # Shortened URL format
        video_id = parsed_url.path.lstrip("/")

    return video_id

def parse_rst(file_path):
    with open(file_path, "r") as f:
        rst_content = f.read()
    # Parse RST into parts (e.g., title, body, etc.)
    parts = publish_parts(rst_content, writer_name="html")
    return parts["body"] 

# Function to get only English transcript and align into 30s segments
def get_english_transcript(video_id, video_name, segment_duration=30):
    try:
        transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
        transcript = transcript_list.find_transcript(['en']).fetch()
        
        # Organize transcript into 30-second segments
        segments = []
        current_segment = []
        start_time = 0

        for entry in transcript:
            if entry["start"] - start_time >= segment_duration:
                segments.append(f"[{format_time(start_time)}] " + " ".join(current_segment))
                start_time = entry["start"]
                current_segment = []
            
            current_segment.append(entry["text"])

        # Save the last segment
        if current_segment:
            segments.append(f"[{format_time(start_time)}] " + " ".join(current_segment))

        # Save to a text file
        file_path = save_transcript(video_name, "\n".join(segments))
        return file_path, "\n".join(segments)

    except Exception as e:
        return None, f"Error: {e}"

def save_pdf_list(pdf_list):
    with open(PDF_LIST_PATH, "w") as f:
        json.dump(pdf_list, f)

def load_pdf_list():
    if os.path.exists(PDF_LIST_PATH):
        with open(PDF_LIST_PATH, "r") as f:
            return json.load(f)
    return []

def analyse_python(content):
    system_prompt = """
    You are a Python learning assistant. Your task is to check whether the query is related to python or not.
    If the query is related to python, you will reply with "Yes". Otherwise, you will reply with "No".
    But only a single word should be a reply.
    Query : {content}
    """

    response = analyser.chat.completions.create(
        model="meta-llama/Llama-3.3-70B-Instruct-Turbo-Free",
        messages=[{"role":"user","content":system_prompt.format(content=content)}]
    )

    return (response.choices[0].message.content)

# Initialize Session State
if "pdf_vector_stores" not in st.session_state:
    st.session_state.pdf_vector_stores = {}

if "error_vector_stores" not in st.session_state:
    st.session_state.error_vector_stores = {}

if "pdf_list" not in st.session_state:
    st.session_state.pdf_list = load_pdf_list()  # Load PDF list from disk

if "messages" not in st.session_state:
    st.session_state.messages = []

if "mode" not in st.session_state:
    st.session_state.mode = "Rubber Duck Assistant"

# Load existing Chroma collections
for pdf_name in st.session_state.pdf_list:
    collection_name = pdf_name.replace(" ", "_").lower()
    vector_store = Chroma(
        collection_name=collection_name,
        embedding_function=EMBEDDING_MODEL,
        persist_directory=CHROMA_DB_PATH
    )
    st.session_state.pdf_vector_stores[pdf_name] = vector_store

for pdf_name in st.session_state.pdf_list:
    collection_name = pdf_name.replace(" ", "_").lower()
    vector_store = Chroma(
        collection_name=collection_name,
        embedding_function=EMBEDDING_MODEL_2,
        persist_directory=CHROMA_DB_PATH_ERROR
    )
    st.session_state.error_vector_stores[pdf_name] = vector_store

def extract_text_from_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()

def extract_text_from_pptx(file_path):
    presentation = Presentation(file_path)
    return "\n".join([shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")])

def embed_query(query):
    return EMBEDDING_MODEL.embed_query(query)

def save_uploaded_file(uploaded_file):
    file_path = os.path.join(PDF_STORAGE_PATH, uploaded_file.name)
    with open(file_path, "wb") as file:
        file.write(uploaded_file.getbuffer())
    return file_path

def process_document(file_name, file_path):
    file_ext = file_name.split(".")[-1].lower()
    
    if file_ext == "pdf":
        loader = PDFPlumberLoader(file_path)
        docs = loader.load()
    elif file_ext == "txt":
        text = extract_text_from_txt(file_path)
        docs = [Document(page_content=text)]
    elif file_ext == "pptx":
        text = extract_text_from_pptx(file_path)
        docs = [Document(page_content=text)]
    elif file_ext == "csv":
        text = extract_text_from_csv(file_path)
        docs = text
    elif file_ext == "rst" or file_ext == "md":
        if file_ext == "rst":
            plain_text = parse_rst(file_path)
            docs = [Document(page_content=plain_text, metadata={"source": file_name})]
        else:
            with open(file_path, "r", encoding="utf-8") as f:
                markdown_content = f.read()
            docs = [Document(page_content=markdown_content, metadata={"source": file_name})]
    else:
        st.error("Unsupported file format.")
        return
    
    if file_ext == "md" or file_ext == "rst":
        chunker = MarkdownTextSplitter(chunk_size=500, chunk_overlap=100)
        chunks = chunker.split_documents(docs)

        documents = [
            Document(
                page_content=chunk.page_content,
                metadata={"source": file_name, "chunk_id": i}
            )
            for i, chunk in enumerate(chunks)
        ]

        collection_name = file_name.replace(" ", "_").lower()

        vector_store = Chroma.from_documents(
            documents=documents,
            collection_name=collection_name,
            persist_directory=CHROMA_DB_PATH_ERROR,
            embedding=EMBEDDING_MODEL_2
        )

        st.session_state.error_vector_stores[file_name] = vector_store
    else:
        # if file_ext == "rst" or file_ext == "md":
        #     chunker = MarkdownTextSplitter(chunk_size=1000, chunk_overlap=200)
        #     chunks = chunker.split_documents(docs)
        # else :
        chunker = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=200)
        chunks = chunker.split_documents(docs)
        
        documents = [
            Document(
                page_content=chunk.page_content,
                metadata={"source": file_name, "chunk_id": i}
            )
            for i, chunk in enumerate(chunks)
        ]
        print("chunks : (FIrst 1000 characters)", documents[:1000])
        collection_name = file_name.replace(" ", "_").lower()
        vector_store = Chroma(
            embedding_function=EMBEDDING_MODEL,
            collection_name=collection_name,
            persist_directory=CHROMA_DB_PATH
        )
        vector_store.add_documents(
            documents=documents,
        )
        st.session_state.pdf_vector_stores[file_name] = vector_store
        stored_data = vector_store.get(include=["embeddings", "documents", "metadatas"])

        # Print embeddings
        for i, embedding in enumerate(stored_data["embeddings"]):
            print(f"Chunk {i} Embedding: {embedding[:5]}... (first 5 values)")
    

    if file_name not in st.session_state.pdf_list:
        st.session_state.pdf_list.append(file_name)
        save_pdf_list(st.session_state.pdf_list)  # Save updated document list to disk

def extract_text_from_csv(file_path):
    import pandas as pd
    
    df = pd.read_csv(file_path)
    text = "\n".join(df.astype(str).apply(lambda row: ", ".join(row), axis=1))  # Convert rows to text
    return text

# Search Relevant Docs
# def find_context(query, selected_pdfs):
#     context_docs = []
#     for pdf in selected_pdfs:
#         store = st.session_state.pdf_vector_stores.get(pdf)
#         if store:
#             context_docs.extend(store.similarity_search(query))
#     return context_docs
def find_context(query, selected_pdfs, mode):
    context_docs = []
    citations = []  # This will store citation objects
    
    for pdf in selected_pdfs:
        # Select vector store (existing code)
        store = (st.session_state.pdf_vector_stores.get(pdf) 
                if mode == "Programming Tutor" 
                else st.session_state.error_vector_stores.get(pdf))

        if not store:
            print(f"⚠️ No vector store found for {pdf}. Skipping.")
            continue
                
        # Perform search (existing code)
        results = store.similarity_search_with_score(query, k=3)
        
        # Enhanced citation handling
        for doc, score in results:
            print("Score:", score, "Doc:", doc)
            if score > 0.7:  # Adjust threshold as needed
                continue
                
            # Extract metadata for citation
            source = doc.metadata.get("source", "Unknown Source")
            page = doc.metadata.get("page", "Unknown")
            chunk_id = doc.metadata.get("chunk_id", "N/A")
            doc_type = doc.metadata.get("type", "content")
            
            # Create unique citation key
            citation_key = f"{source}:{page}:{chunk_id}"
            
            # Store both the document and detailed citation info
            context_docs.append({
                "content": doc.page_content,
                "citation_key": citation_key,
                "score": score
            })
            
            # Build complete citation object
            citations.append({
                "key": citation_key,
                "source": source,
                "page": page,
                "type": doc_type,
                "excerpt": doc.page_content[:200] + "..."  # Preview
            })

    return context_docs, citations

def generate_follow_up_hint(chat_history, mode):
    # Extract the conversation context
    print(chat_history)
    conversation_context = "\n".join([f"{msg[0]}: {msg[1]}" for msg in chat_history])
    conversation_context = conversation_context[:200]
    print(conversation_context)
    # Define the follow-up hint prompt
    if st.session_state.mode == "Programming Tutor":
        prompt = """
        You are a programming tutor. Provide another hint or clarification based on the conversation below.
        Keep your response concise and helpful.

        Conversation:
        {conversation_context}

        Hint:
        """
    elif st.session_state.mode == "Rubber Duck Assistant":
        prompt = """
        You are a rubber duck debugging assistant. Provide another guiding question or hint to help the user think through their problem.

        Conversation:
        {conversation_context}

        Hint:
        """

    # Generate the follow-up hint
    conversation_prompt = ChatPromptTemplate.from_template(prompt)
    response_chain = conversation_prompt | LANGUAGE_MODEL
    return response_chain.invoke({"conversation_context": conversation_context})

def get_chunk_size(query, mode):
    if mode == "Programming Tutor" and len(query) < 50:
        return 200
    elif mode == "Rubber Duck Assistant":
        return 150
    else:
        return 200
# def generate_answer(query, context_docs, citations, mode):
#     chunk_size = get_chunk_size(query, mode)
#     context = "\n\n".join([doc.page_content[:chunk_size] for doc in context_docs])

#     prompt = """
#     You are an AI assistant that provides answers with reliable sources.
#     Use the context below to generate an accurate response.
    
#     Query: {query}
#     Context: {context}
#     Answer:
#     """
    
#     conversation_prompt = ChatPromptTemplate.from_template(prompt)
#     response_chain = conversation_prompt | LANGUAGE_MODEL
#     answer = response_chain.invoke({"query": query, "context": context})

#     # Append citations
#     if citations:
#         sources = "\n".join([f"- {pdf} (Page {page})" for pdf, page in citations])
#         answer += f"\n\n**Sources:**\n{sources}"

#     return answer

def generate_answer(query, context_docs, citations, mode):
    tokenizer = AutoTokenizer.from_pretrained("gpt2")
    prepared_context = []
    used_citation_keys = set()
    current_tokens = 0
    
    for doc in sorted(context_docs, key=lambda x: -x["score"]):
        content_tokens = len(tokenizer.tokenize(doc["content"]))
        if current_tokens + content_tokens > 700:
            continue
        
        prepared_context.append(doc["content"])
        used_citation_keys.add(doc["citation_key"])
        current_tokens += content_tokens

    # chunk_size = get_chunk_size(query, mode)
    # context = "\n\n".join([doc.page_content[:chunk_size] for doc in context_docs])
    print("Context:", prepared_context, "\n")
    # Curriculum-based Assistant Prompt
    if mode == "Programming Tutor":
        prompt = """
        You are a python programming tutor. Explain concepts clearly with examples and best practices.
        Keep explanations concise and beginner-friendly. Use these references etc
        
        Query: {query}
        Context: {context}
        Answer:
        """
    
    # Rubber Duck Debugging Prompt
    elif mode == "Rubber Duck Assistant":
        prompt = """
        You are an expert Rubber Duck Python debugging assistant. Analyze the following Python code and provide debugging hints **one at a time**.
        Do NOT give the full solution immediately.
        Each response should include:
        1️⃣ The next step in debugging.
        2️⃣ A short explanation.
        3️⃣ if the user wants another hint, give it unless it's obvious.
        
        Query: {query}
        Context: {context}
        Answer:
        """

    conversation_prompt = ChatPromptTemplate.from_template(prompt)
    response_chain = conversation_prompt | LANGUAGE_MODEL
    answer = response_chain.invoke({"query": query, "context": prepared_context})

    used_citations = [c for c in citations if c["key"] in used_citation_keys]

    def format_citations(citations):
        """Format citations for display in the UI"""
        formatted = []
        for idx, cite in enumerate(citations, 1):
            formatted.append(
                f"[{idx}] {cite['source']}, page {cite['page']}\n"
                f"Excerpt: {cite['excerpt']}\n"
            )
        return "\n".join(formatted)

    formatted_citations = format_citations(used_citations)
    print("Formatted Citations:",formatted_citations)
    # cited = ""
    # cited_pdfs = set()  # Use a set to store unique citations

    # for pdf in citations:
    #     cited_pdfs.add(f"- {pdf[0]}")  # Add citation as a formatted string

    # # Append citations
    # if cited_pdfs:
    #     sources = "\n".join(sorted(cited_pdfs))  # Sort to maintain order if needed
    #     cited += f"\n\n**Sources:**\n{sources}"

    #     # time.sleep(20)
    return answer,formatted_citations

# Sidebar - Upload PDFs
st.sidebar.header("📤 Upload Programming Resources")
uploaded_files = st.sidebar.file_uploader("Upload PDF, TXT, PPTX, CSV", type=["pdf", "txt", "pptx", "csv", "rst", "md"], accept_multiple_files=True)

# Process New PDFs
if uploaded_files:
    for pdf in uploaded_files:
        if pdf.name not in st.session_state.pdf_vector_stores:
            file_path = save_uploaded_file(pdf)
            process_document(pdf.name, file_path)
    st.sidebar.success("✅ Documents uploaded! Select them below.")

st.title("Programming Tutor "+ "&"+ " Rubber Duck Assistant")
st.markdown("---")

st.sidebar.header("📹 YouTube Transcript Extractor")
st.sidebar.caption("Paste a YouTube link to extract and save the transcript.")

# Input field for YouTube link
video_url = st.sidebar.text_input("🔗 Paste YouTube link here:")

if st.sidebar.button("🎬 Get Transcript"):
    if video_url:
        try:
            # Extract Video ID
            video_id = extract_video_id(video_url)
            if not video_id:
                st.error("❌ Invalid YouTube URL. Please check and try again.")
                st.stop()

            # Generate video name
            video_name = f"transcript_{video_id}"

            # Extract transcript
            file_path, transcript_text = get_english_transcript(video_id, video_name)

            if file_path:
                # Update PDF list
                pdf_list = load_pdf_list()
                if video_name not in pdf_list:
                    pdf_list.append(video_name)
                    save_pdf_list(pdf_list)

                st.sidebar.success(f"✅ Transcript saved as: {video_name}.txt")
                st.rerun()
            else:
                st.error(transcript_text)

        except Exception as e:
            st.sidebar.error(f"❌ Error: {e}")

    else:
        st.warning("⚠ Please enter a valid YouTube link.")

# Sidebar - Select PDFs (Checkbox)
st.sidebar.header("📑 Select Files")
selected_pdfs = [pdf for pdf in st.session_state.pdf_list if st.sidebar.checkbox(pdf, value=True)]

# Sidebar - Delete Files
st.sidebar.header("🗑️ Remove Files")
delete_pdf = st.sidebar.selectbox("Select PDF to Remove", ["None"] + st.session_state.pdf_list)

if st.sidebar.button("❌ Delete Files") and delete_pdf != "None":
    # Delete the Chroma collection for the PDF
    vector_store = st.session_state.pdf_vector_stores.get(delete_pdf)
    if vector_store:
        vector_store.delete_collection()
    del st.session_state.pdf_vector_stores[delete_pdf]
    st.session_state.pdf_list.remove(delete_pdf)
    save_pdf_list(st.session_state.pdf_list)  # Save updated PDF list to disk
    st.sidebar.success(f"Deleted {delete_pdf}")

# Chatbot Mode Selection
mode = st.radio("Choose Assistant Mode:", ["Rubber Duck Assistant","Programming Tutor"])
st.session_state.mode = mode
# Chat Section
if selected_pdfs:
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

    user_query = st.chat_input("Ask a question about programming concepts or debugging...")

    if user_query:
        # Add user query to chat history
        st.session_state.messages.append({"role": "user", "content": user_query})
        with st.chat_message("user"):
            st.markdown(user_query)

        # Generate AI response
        with st.spinner("🔍 Searching resources..."):
            if(analyse_python(user_query) == "Yes"):
                relevant_docs, citations = find_context(user_query, selected_pdfs, st.session_state.mode)
                flag = True
                try:
                    ai_response,cited = generate_answer(user_query, relevant_docs,citations, st.session_state.mode)
                except BadRequestError as e:
                    ai_response = "Sorry, I cannot provide a response to that query due to content filtering policies. Please rephrase your question."
            else:
                flag = False
                ai_response = "Sorry, I cannot provide a response to that query due to content filtering policies. Please rephrase your question."
                st.markdown(ai_response)
                st.session_state.messages.append({"role": "assistant", "content": ai_response})
        # Add AI response to chat history
        if flag:
            with st.chat_message("assistant") :
                st.markdown(ai_response.content)
                st.markdown(cited)
                st.session_state.messages.append({"role": "assistant", "content": ai_response.content})
                st.session_state.messages.append({"role": "assistant", "content": cited})

else:
    st.warning("⚠️ Please upload and select a document to start chatting.")

from streamlit.components.v1 import html
import json
from io import StringIO
import ast
import traceback
from contextlib import redirect_stdout, redirect_stderr


class DebugSandbox:
    def __init__(self):
        self.breakpoints = set()
        self.current_line = 0
        self.variables = {}
        self.execution_log = []
        self.execution_paused = False
        self.breakpoint_hit = False

    def execute_with_debug(self, code, step_mode=False):
        """Execute code with debug capabilities"""
        tree = ast.parse(code)
        self._instrument_code(tree)

        output = StringIO()
        error = StringIO()

        # Add self reference to the execution namespace
        self.variables['self'] = self
        self.variables['__step_mode__'] = step_mode

        try:
            ast.fix_missing_locations(tree)
            with redirect_stdout(output), redirect_stderr(error):
                exec(compile(tree, '<string>', 'exec'), self.variables)
        except Exception:
            error.write(traceback.format_exc())
        finally:
            # Remove self reference to avoid pollution
            if 'self' in self.variables:
                del self.variables['self']
            if '__step_mode__' in self.variables:
                del self.variables['__step_mode__']

        return {
            'output': output.getvalue(),
            'error': error.getvalue(),
            'variables': {k: v for k, v in self.variables.items() if not k.startswith('__')},
            'execution_log': self.execution_log
        }

    def _instrument_code(self, node):
        """Add debugging instrumentation to AST"""
        for field, old_value in ast.iter_fields(node):
            if isinstance(old_value, list):
                new_values = []
                for value in old_value:
                    if isinstance(value, ast.stmt):
                        value = self._instrument_statement(value)
                    new_values.append(value)
                setattr(node, field, new_values)

    def _instrument_statement(self, node):
        self.current_line += 1
        line_num = self.current_line
        
        # Create breakpoint check
        breakpoint_check = ast.If(
            test=ast.Call(
                func=ast.Attribute(
                    value=ast.Name(id='self', ctx=ast.Load()),
                    attr='_check_breakpoint',
                    ctx=ast.Load()
                ),
                args=[ast.Constant(value=line_num)],
                keywords=[]
            ),
            body=[ast.Pass()],
            orelse=[]
        )
        
        # Create step mode check
        step_check = ast.If(
            test=ast.Name(id='__step_mode__', ctx=ast.Load()),
            body=[
                ast.Expr(value=ast.Call(
                    func=ast.Attribute(
                        value=ast.Name(id='self', ctx=ast.Load()),
                        attr='_pause_execution',
                        ctx=ast.Load()
                    ),
                    args=[ast.Constant(value=line_num)],
                    keywords=[]
                ))
            ],
            orelse=[]
        )
        
        # Create log expression
        log_expr = ast.Expr(value=ast.Call(
            func=ast.Attribute(value=ast.Name(id='self', ctx=ast.Load()), 
            attr='_log_execution', 
            ctx=ast.Load()),
            args=[ast.Constant(value=line_num)],
            keywords=[]
        ))

        # Set location info from original node
        log_expr = ast.copy_location(log_expr, node)
        breakpoint_check = ast.copy_location(breakpoint_check, node)
        step_check = ast.copy_location(step_check, node)

        if isinstance(node, ast.FunctionDef):
            node.body = [breakpoint_check, step_check, log_expr] + node.body
            return node
        else:
            new_node = ast.If(
                test=ast.Constant(value=True),
                body=[breakpoint_check, step_check, log_expr, node],
                orelse=[]
            )
            new_node = ast.copy_location(new_node, node)
            return new_node

    def _log_execution(self, line_num):
        """Log execution details"""
        self.execution_log.append({
            'line': line_num,
            'variables': {k: v for k, v in self.variables.items() if not k.startswith('__')},
            'call_stack': traceback.extract_stack()[:-2]
        })

    def _check_breakpoint(self, line_num):
        """Check if execution should pause at this line"""
        if line_num in self.breakpoints:
            self.breakpoint_hit = True
            self.execution_paused = True
            return True
        return False

    def _pause_execution(self, line_num):
        """Pause execution for step-through debugging"""
        self.execution_paused = True
        while self.execution_paused:
            time.sleep(0.1)

    def add_breakpoint(self, line_num):
        """Add a breakpoint at the specified line number"""
        self.breakpoints.add(line_num)

    def remove_breakpoint(self, line_num):
        """Remove breakpoint from the specified line number"""
        if line_num in self.breakpoints:
            self.breakpoints.remove(line_num)

    def continue_execution(self):
        """Resume execution after pause"""
        self.execution_paused = False
        self.breakpoint_hit = False

def generate_answer1(query, code):
    # Enhanced debugging prompt
    prompt = """You are a debugging assistant. Analyze the problem and provide hint-by-hint guidance.
        Context:
        {code}
        Problem: {query}
        Provide:
        1.Give a hint
        2.Error analysis for the error and details about it
        3.Give links/reference to same issues"""

    # Assuming `llm` is a function that takes a prompt and returns a response from a language model
    answer = LANGUAGE_MODEL.invoke(prompt.format(code=code, query=query))

    return answer.content

def get_line_code(full_code, line_num):
    """Returns the specific line of code with line number, plus 2 lines of context"""
    lines = full_code.split('\n')
    start_line = max(1, line_num - 1)  # Show 1 line before (minimum line 1)
    end_line = min(len(lines), line_num + 1)  # Show 1 line after
    
    # Format with line numbers
    result = []
    for i in range(start_line, end_line + 1):
        line_content = lines[i-1]  # Lines are 0-indexed in list
        prefix = ">>" if i == line_num else f"{i:4}"  # Mark error line
        result.append(f"{prefix}: {line_content}")
    
    return '\n'.join(result)

if st.session_state.mode == "Rubber Duck Assistant":
    # Add some CSS for the debug sandbox
    st.markdown("""
        <style>
        .debug-container {
            border: 1px solid #3A3A3A;
            border-radius: 5px;
            padding: 15px;
            margin-bottom: 20px;
            background-color: #1E1E1E;
        }
        .debug-tabs .stTab {
            background-color: #1E1E1E !important;
        }
        </style>
    """, unsafe_allow_html=True)
    
    code = st.text_area("Enter Python code to debug:", height=300, key="debug_code")
    col1, col2 = st.columns([3, 1])


    with col1:
        if st.button("Run with Debugging"):
            sandbox = DebugSandbox()
            result = sandbox.execute_with_debug(st.session_state.debug_code)
            st.session_state.debug_result = result
            st.session_state.stored_code = code
            st.session_state.llm_summary_generated = False
    
    if 'debug_result' in st.session_state:
        result = st.session_state.debug_result
        code = st.session_state.stored_code
        
        st.markdown("### Execution Results")
        tab1, tab3 = st.tabs(["Output", "Errors"])
        
        with tab1:
            st.code(result['output'] or "No output", language='text')
        
        with tab3:
            if result['error']:
                st.error("**Error Details**")
                st.code(result['error'], language='python')
                
                error_lines = result['error'].strip().split('\n')
                if error_lines:
                    last_line = error_lines[-1]
                    if ':' in last_line:
                        error_type, error_message = last_line.split(':', 1)
                        error_type = error_type.strip()
                        error_message = error_message.strip()
                    else:
                        error_type = "Error"
                        error_message = last_line.strip()
                        # answer=generate_answer(error_message,code)
                    st.markdown("**Debug Summary**")
                    st.error(f"**Error Type:** `{error_type}`")
                    st.error(f"**Error Message:** `{error_message}`")
                    error_line = None
                    for line in error_lines:
                        if "line " in line.lower() and ", in " in line.lower():
                            try:
                                line_part = line.split("line ")[1]
                                line_num = line_part.split(",")[0] if "," in line_part else line_part
                                error_line = int(line_num.strip())
                            except (IndexError, ValueError):
                                continue
                    
                    if error_line is not None:
                        st.error(f"**Error Location:** Line {error_line}")
                        st.code(get_line_code(code, error_line), language='python')

                    if st.session_state.llm_summary_generated == False:
                        response_llm = generate_answer1(error_message, code)
                        st.markdown("**Debug Summary from LLM** ")
                        st.markdown(f"{response_llm}")
                        st.session_state.messages.append({"role": "assistant", "content": response_llm})
                        st.session_state.llm_summary_generated = True
                        
                    if st.session_state.mode == "Rubber Duck Assistant" and st.session_state.messages and st.session_state.messages[-1]["role"] == "assistant":
                        if st.button("Need another hint"):
                            # Generate a follow-up hint based on the conversation context
                            with st.spinner("🤔 Generating another hint..."):
                                try:
                                    follow_up_response = generate_follow_up_hint(st.session_state.messages[-1], st.session_state.mode)
                                except BadRequestError as e:
                                    follow_up_response = "Sorry, I cannot provide a response to that query due to content filtering policies. Please rephrase your question."
                                with st.chat_message("assistant"):
                                    st.markdown(follow_up_response.content)
                                    st.session_state.messages.append({"role": "assistant", "content": follow_up_response.content})

            else:
                st.success("No errors detected")